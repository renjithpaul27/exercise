import psycopg2
import pandas as pd
from datetime import datetime
from datetime import date
from airflow.models import Variable
from airflow.operators import PythonOperator
from airflow import DAG
from pptx import Presentation
import itertools
import glob
import os
import re
import csv



sourcefile = '/usr/local/airflow/dags/sourcefile/'+'statusreport.csv'
sourcefileppt = '/usr/local/airflow/dags/sourceppt/'
targetextract = '/usr/local/airflow/dags/sourcefile/'

port = Variable.get("Port")
host = Variable.get("host")
database = Variable.get("database")
username = Variable.get("username")
password = Variable.get("password")

connection = psycopg2.connect(user = username,
                                  password = password,
                                  host = host,
                                  port = port,
                                  database = database)
connection.autocommit = True


def preprocessing(first_page_text):
    first_page_text = re.sub(r'\s+', '  ', first_page_text) # Removing extra spaces
    first_page_text = re.sub(r'\n+', '\n', first_page_text)  # replace multiple newlines with period
    first_page_text = re.sub(r'\.+', '.', first_page_text) # Removing extra periods
    first_page_text = re.sub(r'[^\x00-\x7F]+',' ',first_page_text)# Removing unicode characters
    first_page_text = re.sub(r"\.(?=\S)", ". ", first_page_text)
    return first_page_text

def PPT_process(filename):
    os.chdir(filename)
    filename_csv = targetextract + 'statusreport.csv'
    file_exists = os.path.isfile('filename_csv')
    for eachfile in glob.glob("*.pptx"):
        eachfile_name = os.path.basename(eachfile)
        accountname = eachfile_name.split('_')[0]
        reportingdate = date.today()
        prs = Presentation(eachfile)
        list_of_elements = []
        print("----------------------")
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.is_placeholder:
                    phf = shape.placeholder_format
                    print('%d, %s' % (phf.idx, phf.type))
                if hasattr(shape, "text"):
                    print(shape.text)
                    list_of_elements.append(shape.text)
                if hasattr(shape, "table"):
                    table = shape.table
                    for r in table.rows:
                        s = ""
                        for c in r.cells:
                            s += c.text_frame.text + " | "
                        print(s)
                        list_of_elements.append(s)
                else:
                    pass

        project_title, status, OverallStatus, cost, scope, schedule, resource, risk = data_extraction(list_of_elements)
        accountname_list = [accountname] * len(project_title)
        reportingdate_list = [reportingdate] * len(project_title)

        with open(filename_csv, 'a', newline='') as csvfile:
            writer = csv.writer(csvfile)
            header = ['reportingdate','accountname', 'projectname', 'currentstatus', 'overallstatus', 'coststatus', 'scopestatus', 'schedulestatus', 'resourcingstatus', 'risks']
            if file_exists:
                writer.writerows(zip(reportingdate_list, accountname_list, project_title, status, OverallStatus, cost, scope, schedule, resource, risk))
            else:
                writer.writerow(header)
                file_exists = True
                writer.writerows(zip(reportingdate_list, accountname_list, project_title, status, OverallStatus, cost, scope, schedule, resource, risk))

def data_extraction(list_of_text):
    status = ""
    project_title = []
    current_status = []
    OverallStatus = []
    cost = []
    scope = []
    schedule = []
    resource = []

    risk = []

    for line in range(len(list_of_text)):
        if "Customer Labs" == list_of_text[line].strip():
            project_title.append(list_of_text[line+1])
        if "Current Status" == list_of_text[line].strip():
            status = list_of_text[line+1]
            preprocessed_status = preprocessing(status)
            current_status.append(preprocessed_status)

        if "Overall Status" and 'Scope' in list_of_text[line].strip():
            linev = list_of_text[line+1]
            linev = re.sub(r'\s+', '', linev)
            rag_list = linev.split('|')
            rag_list = list(filter(None, rag_list))
            OverallStatus.append(rag_list[0])
            cost.append(rag_list[1])
            scope.append(rag_list[2])
            schedule.append(rag_list[3])
            resource.append(rag_list[4])
            rag_list = []
        if "CALLOUTS" and "DEPENDENCIES/RISKS" in list_of_text[line].strip():
            risk_value = list_of_text[line]
            risk.append(preprocessing(risk_value))

        else:
            pass
    return project_title, current_status, OverallStatus, cost, scope, schedule, resource, risk

def write_statusreport_file_db(**kwargs):
    '''Reads the source file and laod it to RAW table using copy command'''
    with open(sourcefile, 'r') as statusreportfile:
        cols = statusreportfile.readline().replace("\n","").split(',')
        colsstr = ",".join([str(col) for col in cols])
        print(colsstr)
        with connection.cursor() as cur:
            cur.execute('truncate table raw.statusreportdataraw')
            sql = "copy raw.statusreportdataraw({}) from STDIN WITH delimiter ',' csv".format(colsstr)
            cur.copy_expert(sql,statusreportfile)

def extractppt(**kwargs):
    PPT_process(sourcefileppt)

args = {
    'owner': 'airflow',
    'start_date': datetime(2019, 9, 10, 0, 0),
}

dag = DAG(
    'APAC_STATUS_REPORT_GENERATION',
    schedule_interval="@once",
    default_args=args,
    is_paused_upon_creation=False
)


T1=PythonOperator(
    task_id='Extract_PPT',
    provide_context=True,
    python_callable=extractppt,
    dag=dag,
)


T2=PythonOperator(
    task_id='Write_Source_File',
    provide_context=True,
    python_callable=write_statusreport_file_db,
    dag=dag,
)


T1>>T2
